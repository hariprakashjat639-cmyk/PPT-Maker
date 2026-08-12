import streamlit as st
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
from PIL import Image

# à¤¸à¥à¤°à¤•à¥à¤·à¤¿à¤¤ à¤°à¥‚à¤ª à¤¸à¥‡ à¤µà¥ˆà¤•à¤²à¥à¤ªà¤¿à¤• à¤²à¤¾à¤‡à¤¬à¥à¤°à¥‡à¤°à¥€à¤œ à¤•à¥‹ à¤‡à¤®à¥à¤ªà¥‹à¤°à¥à¤Ÿ à¤•à¤°à¤¨à¤¾ à¤¤à¤¾à¤•à¤¿ à¤à¤ª à¤•à¥à¤°à¥ˆà¤¶ à¤¨ à¤¹à¥‹
try:
    import fitz  # PyMuPDF for PDF & Images
except ImportError:
    fitz = None

try:
    from docx import Document  # python-docx for Word files
except ImportError:
    Document = None

try:
    import pytesseract  # Offline OCR for images
except ImportError:
    pytesseract = None

# à¤ªà¥‡à¤œ à¤¸à¥‡à¤Ÿà¤¿à¤‚à¤—à¥à¤¸ à¤”à¤° à¤²à¥‡à¤†à¤‰à¤Ÿ
st.set_page_config(page_title="Master Offline Model Paper PPT Maker", page_icon="ðŸ“Š", layout="centered")
st.markdown("""
    <style>
    [data-testid="stHeader"] {display: none;}
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    .viewerBadge {display: none !important;}
    header {visibility: hidden;}
    </style>
""", unsafe_allow_html=True)

st.title("ðŸ“Š Master Offline Multi-Format & Direct Text to PPT Maker")
st.write("à¤šà¤¾à¤¹à¥‡ à¤«à¤¾à¤‡à¤² à¤…à¤ªà¤²à¥‹à¤¡ à¤•à¤°à¥‡à¤‚ à¤¯à¤¾ à¤¸à¥€à¤§à¥‡ à¤Ÿà¥‡à¤•à¥à¤¸à¥à¤Ÿ à¤ªà¥‡à¤¸à¥à¤Ÿ à¤•à¤°à¥‡à¤‚à¥¤ à¤¯à¤¹ à¤‘à¤«à¤²à¤¾à¤‡à¤¨ AI à¤”à¤° à¤¡à¤¬à¤²-à¤µà¥‡à¤°à¥€à¤«à¤¾à¤ˆ à¤ªà¤¾à¤°à¥à¤¸à¤° à¤•à¥‡ à¤¸à¤¾à¤¥ à¤†à¤ªà¤•à¥€ à¤ªà¤¸à¤‚à¤¦à¥€à¤¦à¤¾ PPT à¤¤à¥ˆà¤¯à¤¾à¤° à¤•à¤°à¥‡à¤—à¤¾!")

# PPT à¤¸à¥à¤²à¤¾à¤‡à¤¡ à¤¸à¤¾à¤‡à¤œ à¤šà¥à¤¨à¤¨à¥‡ à¤•à¤¾ à¤‘à¤ªà¥à¤¶à¤¨
slide_format = st.selectbox(
    "PPT Slide Size Chunein",
    ["16:9 (Widescreen)", "20:9 (Cinematic)", "4:3 (Standard)"]
)

# à¤‡à¤¨à¤ªà¥à¤Ÿ à¤•à¤¾ à¤¤à¤°à¥€à¤•à¤¾ à¤šà¥à¤¨à¤¨à¥‡ à¤•à¥‡ à¤²à¤¿à¤ à¤µà¤¿à¤•à¤²à¥à¤ª (à¤«à¤¾à¤‡à¤² à¤…à¤ªà¤²à¥‹à¤¡ à¤¯à¤¾ à¤¡à¤¾à¤¯à¤°à¥‡à¤•à¥à¤Ÿ à¤Ÿà¥‡à¤•à¥à¤¸à¥à¤Ÿ à¤ªà¥‡à¤¸à¥à¤Ÿ)
input_choice = st.radio(
    "Data Input Ka Tarika Chunein:",
    ["ðŸ“ File Upload Karein (.txt, .pdf, .docx, Image)", "âœï¸ Direct Text Paste Karein"]
)

raw_text = ""

if input_choice == "ðŸ“ File Upload Karein (.txt, .pdf, .docx, Image)":
    uploaded_file = st.file_uploader(
        "File Upload Karein (.txt, .pdf, .docx, .png, .jpg, .jpeg)", 
        type=["txt", "pdf", "docx", "png", "jpg", "jpeg"]
    )
    if uploaded_file is not None:
        file_bytes = uploaded_file.getvalue()
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        # --- à¤µà¤¿à¤­à¤¿à¤¨à¥à¤¨ à¤«à¥‰à¤°à¥à¤®à¥‡à¤Ÿà¥à¤¸ à¤¸à¥‡ à¤Ÿà¥‡à¤•à¥à¤¸à¥à¤Ÿ à¤¨à¤¿à¤•à¤¾à¤²à¤¨à¥‡ à¤µà¤¾à¤²à¤¾ à¤«à¤‚à¤•à¥à¤¶à¤¨ ---
        try:
            if file_extension == "txt":
                raw_text = file_bytes.decode("utf-8", errors="ignore")
            elif file_extension == "pdf":
                if fitz is None:
                    st.error("âš ï¸ PyMuPDF (fitz) à¤²à¤¾à¤‡à¤¬à¥à¤°à¥‡à¤°à¥€ à¤‡à¤‚à¤¸à¥à¤Ÿà¥‰à¤² à¤¨à¤¹à¥€à¤‚ à¤¹à¥ˆà¥¤")
                else:
                    doc = fitz.open(stream=file_bytes, filetype="pdf")
                    for page in doc:
                        raw_text += page.get_text() + "\n"
            elif file_extension == "docx":
                if Document is None:
                    st.error("âš ï¸ python-docx à¤²à¤¾à¤‡à¤¬à¥à¤°à¥‡à¤°à¥€ à¤‡à¤‚à¤¸à¥à¤Ÿà¥‰à¤² à¤¨à¤¹à¥€à¤‚ à¤¹à¥ˆà¥¤")
                else:
                    doc = Document(io.BytesIO(file_bytes))
                    for para in doc.paragraphs:
                        raw_text += para.text + "\n"
            elif file_extension in ["png", "jpg", "jpeg"]:
                if pytesseract is None:
                    st.error("âš ï¸ pytesseract à¤²à¤¾à¤‡à¤¬à¥à¤°à¥‡à¤°à¥€ à¤‡à¤‚à¤¸à¥à¤Ÿà¥‰à¤² à¤¨à¤¹à¥€à¤‚ à¤¹à¥ˆà¥¤")
                else:
                    image = Image.open(io.BytesIO(file_bytes))
                    raw_text = pytesseract.image_to_string(image, lang='hin+eng')
        except Exception as e:
            st.error(f"File reading error: {e}")
else:
    raw_text = st.text_area(
        "Yahan Apne Prashn, Options, Answer aur Explanation Paste Karein:",
        height=250,
        placeholder="à¤ªà¥à¤°à¤¶à¥à¤¨ 1: à¤­à¤¾à¤°à¤¤ à¤•à¥€ à¤°à¤¾à¤œà¤§à¤¾à¤¨à¥€ à¤•à¥à¤¯à¤¾ à¤¹à¥ˆ?\n(A) à¤®à¥à¤‚à¤¬à¤ˆ\n(B) à¤¦à¤¿à¤²à¥à¤²à¥€\n(C) à¤•à¥‹à¤²à¤•à¤¾à¤¤à¤¾\n(D) à¤šà¥‡à¤¨à¥à¤¨à¤ˆ\nà¤‰à¤¤à¥à¤¤à¤°: (B) à¤¦à¤¿à¤²à¥à¤²à¥€\nà¤µà¥à¤¯à¤¾à¤–à¥à¤¯à¤¾: à¤¦à¤¿à¤²à¥à¤²à¥€ à¤­à¤¾à¤°à¤¤ à¤•à¥€ à¤°à¤¾à¤·à¥à¤Ÿà¥à¤°à¥€à¤¯ à¤°à¤¾à¤œà¤§à¤¾à¤¨à¥€ à¤¹à¥ˆà¥¤"
    )

# --- à¤¡à¤¬à¤²-à¤µà¥‡à¤°à¥€à¤«à¤¾à¤ˆ à¤”à¤° à¤¸à¥à¤®à¤¾à¤°à¥à¤Ÿ à¤Ÿà¥‡à¤•à¥à¤¸à¥à¤Ÿ à¤ªà¤¾à¤°à¥à¤¸à¤° ---
def double_verify_and_parse(text):
    questions = []
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    
    current_q = None
    q_pattern = re.compile(r'^(à¤ªà¥à¤°à¤¶à¥à¤¨|\bQ\b|\bQ\d+|\d+[\.\)]|\bQuestion\b)', re.IGNORECASE)
    opt_pattern = re.compile(r'^(\([A-Da-dà¤…-à¤¦1-4]\)|[A-Da-dà¤…-à¤¦1-4][\.\)]|\b[A-Da-dà¤…-à¤¦][\)])')
    ans_pattern = re.compile(r'^(à¤‰à¤¤à¥à¤¤à¤°|Answer|Ans|à¤¸à¤¹à¥€ à¤‰à¤¤à¥à¤¤à¤°)[\s:\-]', re.IGNORECASE)
    exp_pattern = re.compile(r'^(à¤µà¥à¤¯à¤¾à¤–à¥à¤¯à¤¾|Explanation|Exp|à¤¸à¥à¤ªà¤·à¥à¤Ÿà¥€à¤•à¤°à¤£)[\s:\-]', re.IGNORECASE)

    for line in lines:
        if q_pattern.match(line) or (current_q and line.startswith('à¤ªà¥à¤°à¤¶à¥à¤¨')):
            if current_q and current_q['question']:
                while len(current_q['options']) < 4:
                    current_q['options'].append(f"({len(current_q['options']) + 1}) à¤µà¤¿à¤•à¤²à¥à¤ª à¤‰à¤ªà¤²à¤¬à¥à¤§ à¤¨à¤¹à¥€à¤‚")
                questions.append(current_q)
            current_q = {'question': line, 'options': [], 'answer': '', 'explanation': []}
            continue
            
        if current_q is None:
            current_q = {'question': line, 'options': [], 'answer': '', 'explanation': []}
            continue

        if opt_pattern.match(line):
            current_q['options'].append(line)
        elif ans_pattern.match(line):
            current_q['answer'] = line
        elif exp_pattern.match(line):
            exp_text = exp_pattern.sub('', line).strip()
            if exp_text:
                current_q['explanation'].append(exp_text)
        else:
            if current_q['explanation']:
                if len(current_q['explanation']) < 3:
                    current_q['explanation'].append(line)
            elif current_q['answer']:
                current_q['answer'] += " " + line
            elif current_q['options']:
                if len(current_q['options']) < 4:
                    current_q['options'].append(line)
                else:
                    current_q['explanation'].append(line)
            else:
                current_q['question'] += " " + line

    if current_q and current_q['question']:
        while len(current_q['options']) < 4:
            current_q['options'].append(f"({len(current_q['options']) + 1}) à¤µà¤¿à¤•à¤²à¥à¤ª à¤‰à¤ªà¤²à¤¬à¥à¤§ à¤¨à¤¹à¥€à¤‚")
        questions.append(current_q)
        
    return questions

if st.button("ðŸš€ Master PPT Generate Karein"):
    if not raw_text.strip():
        st.warning("âš ï¸ à¤•à¥ƒà¤ªà¤¯à¤¾ à¤ªà¤¹à¤²à¥‡ à¤•à¥‹à¤ˆ à¤«à¤¾à¤‡à¤² à¤…à¤ªà¤²à¥‹à¤¡ à¤•à¤°à¥‡à¤‚ à¤¯à¤¾ à¤Ÿà¥‡à¤•à¥à¤¸à¥à¤Ÿ à¤¬à¥‰à¤•à¥à¤¸ à¤®à¥‡à¤‚ à¤ªà¥à¤°à¤¶à¥à¤¨/à¤‰à¤¤à¥à¤¤à¤° à¤ªà¥‡à¤¸à¥à¤Ÿ à¤•à¤°à¥‡à¤‚!")
    else:
        with st.spinner("à¤Ÿà¥‡à¤•à¥à¤¸à¥à¤Ÿ à¤•à¥‹ à¤ªà¤¾à¤°à¥à¤¸ à¤•à¤¿à¤¯à¤¾ à¤œà¤¾ à¤°à¤¹à¤¾ à¤¹à¥ˆ à¤”à¤° à¤¡à¤¬à¤²-à¤µà¥‡à¤°à¥€à¤«à¤¾à¤ˆ à¤•à¤¿à¤¯à¤¾ à¤œà¤¾ à¤°à¤¹à¤¾ à¤¹à¥ˆ..."):
            parsed_questions = double_verify_and_parse(raw_text)
            
            prs = Presentation()
            
            # à¤¤à¥€à¤¨à¥‹à¤‚ à¤¸à¤¾à¤‡à¤œà¤¼ à¤”à¤° à¤‰à¤¨à¤•à¥€ à¤ªà¤°à¤«à¥‡à¤•à¥à¤Ÿ à¤¡à¤¿à¤œà¤¼à¤¾à¤‡à¤¨ à¤¸à¥‡à¤Ÿà¤¿à¤‚à¤—à¥à¤¸
            if slide_format == "20:9 (Cinematic)":
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(6.0)
                q_font_size = Pt(32)
                opt_font_size = Pt(30)
                ans_font_size = Pt(23)
                exp_font_size = Pt(30)
                card_width = Inches(12.333)
                card_height = Inches(5.0)
                box_width = Inches(11.733)
                q_box_width = Inches(12.3)
                opt_box_width = Inches(12.0)
                exp_box_height = Inches(3.2)
                opt_left = Inches(0.9)
                opt_top = Inches(2.5)
                opt_space_before = Pt(2)

            elif slide_format == "16:9 (Widescreen)":
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                q_font_size = Pt(40)
                opt_font_size = Pt(40)
                ans_font_size = Pt(28)
                exp_font_size = Pt(32)
                card_width = Inches(11.733)
                card_height = Inches(6.4)
                box_width = Inches(11.133)
                q_box_width = Inches(12.3)
                opt_box_width = Inches(12.0)
                exp_box_height = Inches(4.5)
                opt_left = Inches(0.8)
                opt_top = Inches(3.0)
                opt_space_before = Pt(8)

            elif slide_format == "4:3 (Standard)":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                q_font_size = Pt(30)
                opt_font_size = Pt(28)
                ans_font_size = Pt(24)
                exp_font_size = Pt(24)
                card_width = Inches(8.8)
                card_height = Inches(6.4)
                box_width = Inches(8.2)
                q_box_width = Inches(9.0)
                opt_box_width = Inches(8.8)
                exp_box_height = Inches(4.5)
                opt_left = Inches(0.5)
                opt_top = Inches(2.8)
                opt_space_before = Pt(6)

            blank_layout = prs.slide_layouts[6] 

            for idx, q in enumerate(parsed_questions):
                # ==========================================
                # SLIDE 1: Question + Options
                # ==========================================
                slide1 = prs.slides.add_slide(blank_layout)

                q_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.4), q_box_width, Inches(2.5))
                tf1 = q_box.text_frame
                tf1.word_wrap = True
                p1 = tf1.paragraphs[0]
                p1.text = q['question']
                p1.font.name = 'Nirmala UI'
                p1.font.size = q_font_size
                p1.font.bold = True
                p1.font.color.rgb = RGBColor(255, 0, 0)  # Pure Red (#FF0000)
                p1.line_spacing = 1.3

                opt_box = slide1.shapes.add_textbox(opt_left, opt_top, opt_box_width, Inches(4.3))
                tf_opt = opt_box.text_frame
                tf_opt.word_wrap = True

                options = q['options'] if q['options'] else ["(A) à¤µà¤¿à¤•à¤²à¥à¤ª 1", "(B) à¤µà¤¿à¤•à¤²à¥à¤ª 2", "(C) à¤µà¤¿à¤•à¤²à¥à¤ª 3", "(D) à¤µà¤¿à¤•à¤²à¥à¤ª 4"]
                for opt_idx, opt in enumerate(options):
                    if opt_idx == 0:
                        p = tf_opt.paragraphs[0]
                    else:
                        p = tf_opt.add_paragraph()
                        p.space_before = opt_space_before
                    
                    p.text = opt
                    p.font.name = 'Nirmala UI'
                    p.font.size = opt_font_size
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.line_spacing = 1.4

                # ==========================================
                # SLIDE 2: Answer + Explanation
                # ==========================================
                slide2 = prs.slides.add_slide(blank_layout)

                card = slide2.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    Inches(0.8), Inches(0.5), card_width, card_height
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(248, 250, 252)
                card.line.color.rgb = RGBColor(203, 213, 225)
                card.line.width = Pt(1.5)

                ans_banner = slide2.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1.1), Inches(0.8), box_width, Inches(1.1)
                )
                ans_banner.fill.solid()
                ans_banner.fill.fore_color.rgb = RGBColor(22, 163, 74) # Green Banner
                ans_banner.line.color.rgb = RGBColor(22, 163, 74)

                tf_ans = ans_banner.text_frame
                tf_ans.word_wrap = True
                p_ans = tf_ans.paragraphs[0]
                p_ans.text = q['answer'] if q['answer'] else "à¤‰à¤¤à¥à¤¤à¤°: (à¤¸à¤¹à¥€ à¤µà¤¿à¤•à¤²à¥à¤ª à¤•à¤¾ à¤¨à¤¾à¤®)"
                p_ans.font.name = 'Nirmala UI'
                p_ans.font.size = ans_font_size
                p_ans.font.bold = True
                p_ans.font.color.rgb = RGBColor(255, 255, 255)

                exp_box = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), box_width, exp_box_height)
                tf_exp = exp_box.text_frame
                tf_exp.word_wrap = True

                p_exp_title = tf_exp.paragraphs[0]
                p_exp_title.text = "à¤µà¥à¤¯à¤¾à¤–à¥à¤¯à¤¾:"
                p_exp_title.font.name = 'Nirmala UI'
                p_exp_title.font.size = Pt(26)
                p_exp_title.font.bold = True
                p_exp_title.font.color.rgb = RGBColor(30, 41, 59)

                expl_lines = q['explanation'][:3] if q['explanation'] else ["à¤®à¤¹à¤¤à¥à¤µà¤ªà¥‚à¤°à¥à¤£ à¤µà¥à¤¯à¤¾à¤–à¥à¤¯à¤¾ à¤¬à¤¿à¤‚à¤¦à¥ à¤¯à¤¹à¤¾à¤ à¤†à¤à¤‚à¤—à¥‡à¥¤"]

                for line_idx, exp_line in enumerate(expl_lines):
                    p_exp = tf_exp.add_paragraph()
                    p_exp.space_before = Pt(6)
                    p_exp.line_spacing = 1.25
                    
                    p_exp.text = f"â€¢ {exp_line}" if not exp_line.startswith('â€¢') else exp_line
                    p_exp.font.name = 'Nirmala UI'
                    p_exp.font.size = exp_font_size
                    p_exp.font.color.rgb = RGBColor(0, 0, 0)

            ppt_stream = io.BytesIO()
            prs.save(ppt_stream)
            ppt_stream.seek(0)

            st.success("ðŸŽ‰ à¤†à¤ªà¤•à¥€ à¤®à¤¾à¤¸à¥à¤Ÿà¤° PPT à¤¸à¤«à¤²à¤¤à¤¾àªªà¥‚à¤°à¥à¤µà¤• à¤¤à¥ˆà¤¯à¤¾à¤° à¤¹à¥‹ à¤—à¤ˆ à¤¹à¥ˆ!")
            st.download_button(
                label="ðŸ“¥ PPT Download Karein",
                data=ppt_stream,
                file_name="Master_Model_Paper.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
